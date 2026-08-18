#!/usr/bin/env python3
"""Génère la présentation PowerPoint de la plateforme IoT LoRaWAN."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

ORANGE = RGBColor(0xFF, 0x79, 0x00)
DARK = RGBColor(0x1A, 0x1A, 0x2E)
GRAY = RGBColor(0x55, 0x55, 0x55)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_ORANGE = RGBColor(0xFF, 0xF3, 0xE6)
GREEN = RGBColor(0x1B, 0x7D, 0x3A)
RED = RGBColor(0xC0, 0x39, 0x2B)
AMBER = RGBColor(0xE6, 0x7E, 0x22)
OUTPUT = Path(__file__).resolve().parent / "Plateforme_IoT_LoRaWAN_Sonatel.pptx"


def set_slide_bg(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title_slide(prs: Presentation, title: str, subtitle: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK)
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.15), prs.slide_height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = ORANGE
    bar.line.fill.background()

    box = slide.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(11), Inches(1.5))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE

    sub = slide.shapes.add_textbox(Inches(0.8), Inches(3.8), Inches(11), Inches(1.2))
    sp = sub.text_frame.paragraphs[0]
    sp.text = subtitle
    sp.font.size = Pt(20)
    sp.font.color.rgb = ORANGE

    foot = slide.shapes.add_textbox(Inches(0.8), Inches(6.5), Inches(11), Inches(0.5))
    fp = foot.text_frame.paragraphs[0]
    fp.text = "Sonatel · Orange IoT  |  LoRaWAN SaaS Platform  |  2026"
    fp.font.size = Pt(14)
    fp.font.color.rgb = GRAY


def add_section_slide(prs: Presentation, title: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, ORANGE)
    box = slide.shapes.add_textbox(Inches(0.8), Inches(2.8), Inches(11), Inches(1.5))
    p = box.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT


def add_content_slide(
    prs: Presentation,
    title: str,
    bullets: list[str],
    subtitle: str | None = None,
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    header = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, Inches(1.1))
    header.fill.solid()
    header.fill.fore_color.rgb = DARK
    header.line.fill.background()

    accent = slide.shapes.add_shape(1, Inches(0), Inches(1.1), prs.slide_width, Inches(0.06))
    accent.fill.solid()
    accent.fill.fore_color.rgb = ORANGE
    accent.line.fill.background()

    tb = slide.shapes.add_textbox(Inches(0.6), Inches(0.25), Inches(12), Inches(0.7))
    tp = tb.text_frame.paragraphs[0]
    tp.text = title
    tp.font.size = Pt(28)
    tp.font.bold = True
    tp.font.color.rgb = WHITE

    top = Inches(1.5)
    if subtitle:
        st = slide.shapes.add_textbox(Inches(0.7), Inches(1.35), Inches(12), Inches(0.5))
        sp = st.text_frame.paragraphs[0]
        sp.text = subtitle
        sp.font.size = Pt(14)
        sp.font.italic = True
        sp.font.color.rgb = ORANGE
        top = Inches(1.85)

    body = slide.shapes.add_textbox(Inches(0.7), top, Inches(12), Inches(5.5))
    tf = body.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(17)
        p.font.color.rgb = DARK
        p.space_after = Pt(10)


def add_two_column_slide(
    prs: Presentation,
    title: str,
    left_title: str,
    left_bullets: list[str],
    right_title: str,
    right_bullets: list[str],
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    header = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, Inches(1.1))
    header.fill.solid()
    header.fill.fore_color.rgb = DARK
    header.line.fill.background()

    tb = slide.shapes.add_textbox(Inches(0.6), Inches(0.25), Inches(12), Inches(0.7))
    tp = tb.text_frame.paragraphs[0]
    tp.text = title
    tp.font.size = Pt(28)
    tp.font.bold = True
    tp.font.color.rgb = WHITE

    for col, col_title, bullets, x in [
        (0, left_title, left_bullets, 0.5),
        (1, right_title, right_bullets, 6.8),
    ]:
        card = slide.shapes.add_shape(1, Inches(x), Inches(1.4), Inches(6.0), Inches(5.8))
        card.fill.solid()
        card.fill.fore_color.rgb = LIGHT_ORANGE
        card.line.color.rgb = ORANGE

        ct = slide.shapes.add_textbox(Inches(x + 0.2), Inches(1.55), Inches(5.6), Inches(0.5))
        cp = ct.text_frame.paragraphs[0]
        cp.text = col_title
        cp.font.size = Pt(18)
        cp.font.bold = True
        cp.font.color.rgb = ORANGE

        body = slide.shapes.add_textbox(Inches(x + 0.2), Inches(2.1), Inches(5.6), Inches(4.8))
        tf = body.text_frame
        tf.word_wrap = True
        for i, bullet in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = bullet
            p.font.size = Pt(14)
            p.font.color.rgb = DARK
            p.space_after = Pt(6)


def _slide_header(slide, prs: Presentation, title: str, subtitle: str | None = None) -> None:
    set_slide_bg(slide, WHITE)
    header = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, Inches(1.1))
    header.fill.solid()
    header.fill.fore_color.rgb = DARK
    header.line.fill.background()
    accent = slide.shapes.add_shape(1, Inches(0), Inches(1.1), prs.slide_width, Inches(0.06))
    accent.fill.solid()
    accent.fill.fore_color.rgb = ORANGE
    accent.line.fill.background()
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(0.25), Inches(12), Inches(0.7))
    tp = tb.text_frame.paragraphs[0]
    tp.text = title
    tp.font.size = Pt(26)
    tp.font.bold = True
    tp.font.color.rgb = WHITE
    if subtitle:
        st = slide.shapes.add_textbox(Inches(0.6), Inches(0.72), Inches(12), Inches(0.35))
        sp = st.text_frame.paragraphs[0]
        sp.text = subtitle
        sp.font.size = Pt(11)
        sp.font.italic = True
        sp.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)


def add_compliance_slide(prs: Presentation) -> None:
    """Comparatif bonnes pratiques Actility/ThingPark vs plateforme Sonatel."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_header(
        slide,
        prs,
        "Conformité vs Actility / ThingPark",
        "Référentiel industriel LNS — état actuel de la plateforme Sonatel",
    )

    rows = [
        ("LNS redondant multi-sites (2–3)", "Requis industriel", "Mono-VM Docker Compose", "Non", RED),
        ("VPN gateway ↔ LNS + PKI", "Obligatoire", "UDP :1700 ouvert en prod", "Non", RED),
        ("HSM / clés FIPS 140", "Recommandé utilities", "Clés ChirpStack, pas de HSM", "Non", RED),
        ("Routage + decoders par route", "Connectors ThingPark", "Connecteurs HTTP/MQTT/MCP + codecs JS", "Oui", GREEN),
        ("Connecteurs HA primaire/secondaire", "Requis", "Un connecteur par tenant", "Partiel", AMBER),
        ("ADR (Adaptive Data Rate)", "Recommandé", "ChirpStack ADR activé (DR5/SF7)", "Oui", GREEN),
        ("Supervision NOC + maintenance", "Requis", "NOC, anomalies, RF Scan, agent IA", "Oui", GREEN),
        ("FUOTA + commissioning ICS", "Requis à l'échelle", "FUOTA Phase 3, pas d'ICS centralisé", "Partiel", AMBER),
        ("Peering / NetID LoRa Alliance", "Réseau national", "Non configuré", "Non", RED),
        ("Migration export/import clés", "Labo → national", "Possible ChirpStack, non automatisé", "Partiel", AMBER),
        ("Reprise progressive (packet storm)", "Industriel", "Non implémenté", "Non", RED),
    ]

    table_shape = slide.shapes.add_table(
        len(rows) + 1, 4, Inches(0.35), Inches(1.25), Inches(12.6), Inches(5.35)
    )
    table = table_shape.table
    col_widths = [Inches(2.7), Inches(2.2), Inches(3.5), Inches(1.2)]
    for i, w in enumerate(col_widths):
        table.columns[i].width = w

    headers = ["Domaine", "Exigence Actility", "Plateforme Sonatel", "Statut"]
    for col, text in enumerate(headers):
        cell = table.cell(0, col)
        cell.text = text
        cell.fill.solid()
        cell.fill.fore_color.rgb = DARK
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(10)
            p.font.bold = True
            p.font.color.rgb = WHITE

    for row_idx, (domain, actility, sonatel, status, color) in enumerate(rows, start=1):
        values = [domain, actility, sonatel, status]
        for col, text in enumerate(values):
            cell = table.cell(row_idx, col)
            cell.text = text
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(9)
                p.font.color.rgb = DARK
            if col == 3:
                for p in cell.text_frame.paragraphs:
                    p.font.bold = True
                    p.font.color.rgb = color
            if row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xFA, 0xFA, 0xFA)

    foot = slide.shapes.add_textbox(Inches(0.35), Inches(6.65), Inches(12.6), Inches(0.55))
    fp = foot.text_frame.paragraphs[0]
    fp.text = (
        "Verdict : prêt pour pilote (centaines de devices) · gaps infrastructure pour déploiement national "
        "(VPN, HA, HSM, peering) — Phase 4 roadmap"
    )
    fp.font.size = Pt(10)
    fp.font.bold = True
    fp.font.color.rgb = ORANGE


def add_architecture_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    header = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, Inches(1.1))
    header.fill.solid()
    header.fill.fore_color.rgb = DARK
    header.line.fill.background()

    tb = slide.shapes.add_textbox(Inches(0.6), Inches(0.25), Inches(12), Inches(0.7))
    tp = tb.text_frame.paragraphs[0]
    tp.text = "Architecture de la plateforme"
    tp.font.size = Pt(28)
    tp.font.bold = True
    tp.font.color.rgb = WHITE

    layers = [
        ("Console Web (Next.js)", "Portail Orange Live Objects — devices, NOC, analytics, agent IA"),
        ("Platform API (Go)", "Multi-tenant, proxy ChirpStack, billing, règles, apps métier"),
        ("ChirpStack v4", "Network Server LoRaWAN — NS, Join Server, Application Server"),
        ("Ingestion & Data", "MQTT → TimescaleDB + MinIO + NATS (rules, connecteurs)"),
        ("Agent IA (MCP)", "~25 outils — provisioning, diagnostics, compteurs eau en langage naturel"),
    ]
    y = 1.5
    for i, (name, desc) in enumerate(layers):
        color = ORANGE if i == 0 else DARK
        box = slide.shapes.add_shape(1, Inches(1.5), Inches(y), Inches(10), Inches(0.85))
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.fill.background()

        nt = slide.shapes.add_textbox(Inches(1.7), Inches(y + 0.08), Inches(9.6), Inches(0.35))
        np = nt.text_frame.paragraphs[0]
        np.text = name
        np.font.size = Pt(16)
        np.font.bold = True
        np.font.color.rgb = WHITE

        dt = slide.shapes.add_textbox(Inches(1.7), Inches(y + 0.42), Inches(9.6), Inches(0.35))
        dp = dt.text_frame.paragraphs[0]
        dp.text = desc
        dp.font.size = Pt(12)
        dp.font.color.rgb = WHITE if i == 0 else RGBColor(0xCC, 0xCC, 0xCC)

        if i < len(layers) - 1:
            arrow = slide.shapes.add_textbox(Inches(6.2), Inches(y + 0.85), Inches(0.5), Inches(0.3))
            ap = arrow.text_frame.paragraphs[0]
            ap.text = "▼"
            ap.font.size = Pt(14)
            ap.font.color.rgb = ORANGE
            ap.alignment = PP_ALIGN.CENTER

        y += 1.05

    side = slide.shapes.add_textbox(Inches(0.5), Inches(6.2), Inches(12), Inches(0.8))
    sp = side.text_frame.paragraphs[0]
    sp.text = (
        "Stack : Keycloak (IAM) · Redis · NATS · MinIO · Stripe · "
        "Déploiement Docker — EU868, OTAA, Class A/B/C, FUOTA"
    )
    sp.font.size = Pt(11)
    sp.font.color.rgb = GRAY


def build() -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    add_title_slide(
        prs,
        "Plateforme IoT LoRaWAN Nationale",
        "Souveraineté numérique · Smart metering · Opérations réseau · Intelligence artificielle",
    )

    add_content_slide(
        prs,
        "Sommaire",
        [
            "1. Contexte et opportunités marché",
            "2. Enjeux stratégiques",
            "3. Objectifs de la plateforme",
            "4. Présentation de la solution",
            "5. Fonctionnalités clés",
            "6. Applications métier (eau intelligente)",
            "7. Modèle économique et différenciation",
            "8. Roadmap et prochaines étapes",
        ],
    )

    add_section_slide(prs, "Contexte & opportunités")

    add_content_slide(
        prs,
        "Un marché IoT en forte accélération",
        [
            "Explosion des cas d'usage LPWAN : comptage eau/énergie, agriculture, smart city, industrie",
            "LoRaWAN : technologie ouverte, longue portée, faible consommation — idéale pour le Sénégal",
            "Besoin croissant de connecter des millions de capteurs à faible coût (CAPEX et OPEX maîtrisés)",
            "Digitalisation des utilities (eau, électricité) : télérelevé, vanne intelligente, détection de fuites",
            "Pression réglementaire et environnementale : réduction des pertes en eau (NRW — Non-Revenue Water)",
            "Opportunité pour un opérateur télécom de devenir acteur IoT vertical, au-delà de la connectivité",
        ],
        subtitle="Afrique de l'Ouest · Sénégal · 2025–2030",
    )

    add_content_slide(
        prs,
        "SUEZ — Vers un réseau LoRaWAN national eau",
        [
            "SUEZ (leader mondial de l'eau) engage des discussions pour un déploiement LoRaWAN à l'échelle nationale",
            "Objectif : télérelevé massif de compteurs intelligents, supervision des réseaux d'eau potable",
            "Cas d'usage : index de consommation, commande de vannes à distance, détection précoce de fuites",
            "Enjeu : couvrir les zones urbaines et périurbaines (Dakar, régions) avec un réseau unifié",
            "Sonatel / Orange : position unique — infrastructure télécom + couverture nationale + data centers",
            "La plateforme développée répond directement à ce besoin avec une solution eau prête (Shengda V1.6)",
        ],
    )

    add_two_column_slide(
        prs,
        "Coût des licences LNS propriétaires",
        "Actility ThingPark",
        [
            "LNS cloud ou on-premise — référence marché",
            "Licence par gateway et/ou par device",
            "Coûts récurrents élevés à l'échelle nationale",
            "Vendor lock-in — dépendance fournisseur",
            "Personnalisation et apps métier limitées",
            "Facturation opaque pour grands déploiements",
        ],
        "Netmore · Loriot · autres",
        [
            "Modèles SaaS avec quotas devices/uplinks",
            "Coût prohibitif au-delà de quelques milliers de devices",
            "Peu adaptés au contexte africain (support, latence)",
            "Pas de souveraineté des données sur le territoire",
            "Intégration SI client complexe et coûteuse",
            "Pas d'agent IA opérationnel natif",
        ],
    )

    add_content_slide(
        prs,
        "Programme New Deal — Souveraineté IoT nationale",
        [
            "L'État du Sénégal, via le programme New Deal Technologique, vise une plateforme IoT nationale souveraine",
            "Objectif : réduire la dépendance aux solutions étrangères et maîtriser les données stratégiques",
            "Alignement avec la stratégie numérique du Sénégal (SENEGEL 2025, smart cities, e-gov)",
            "Interopérabilité avec les opérateurs, utilities et collectivités locales",
            "Hébergement local des données (conformité, souveraineté, latence)",
            "Sonatel comme partenaire naturel : réseau, expertise télécom, ancrage local",
        ],
    )

    add_content_slide(
        prs,
        "Autres éléments de contexte",
        [
            "SENELEC / comptage énergie : même logique LPWAN que l'eau (télérelevé, supervision)",
            "Smart cities : gestion des déchets, éclairage public, parking, qualité de l'air",
            "Agriculture : capteurs sol/humidité, irrigation intelligente (rural LoRaWAN)",
            "Industrie & logistique : suivi actifs, chaîne du froid, maintenance prédictive",
            "Région EU868 déjà opérationnelle — gateways déployées, devices OTAA en production",
            "ChirpStack open-source (v4) : base éprouvée, évite 12–18 mois de R&D Network Server",
        ],
    )

    add_section_slide(prs, "Enjeux & objectifs")

    add_two_column_slide(
        prs,
        "Enjeux stratégiques",
        "Économiques",
        [
            "Éviter les licences LNS à coût prohibitif (Actility, Netmore…)",
            "Modèle SaaS récurrent vs achat de licence par device",
            "Time-to-market : déploiement en mois, pas en années",
            "Réduction des pertes en eau = ROI direct pour les utilities",
            "Nouvelle ligne de revenus IoT pour Sonatel",
        ],
        "Techniques & institutionnels",
        [
            "Souveraineté des données et hébergement local",
            "Multi-tenant : un réseau, plusieurs clients (SUEZ, villes, industriels)",
            "Scalabilité : de 100 à 1 M+ devices",
            "Sécurité : IAM Keycloak, isolation tenant, API keys",
            "Interopérabilité : webhooks, MQTT, MCP vers ERP/SCADA clients",
        ],
    )

    add_content_slide(
        prs,
        "Objectifs de la plateforme",
        [
            "Construire une plateforme IoT LoRaWAN SaaS cloud-native, propriété Sonatel / Orange",
            "Offrir une alternative souveraine et économique aux LNS propriétaires (Actility, Netmore…)",
            "Supporter un déploiement national de compteurs eau intelligents (SUEZ et autres utilities)",
            "Répondre aux exigences du programme New Deal (plateforme IoT nationale)",
            "Industrialiser les opérations réseau : NOC, anomalies, diagnostics IA",
            "Accélérer l'onboarding client : wizard, agent IA, provisioning automatique multi-tenant",
            "Monétiser via plans Starter / Pro / Enterprise avec metering et facturation intégrée",
        ],
    )

    add_section_slide(prs, "La plateforme LoRaWAN SaaS")

    add_content_slide(
        prs,
        "Proposition de valeur",
        [
            "Plateforme IoT LoRaWAN complète — de la gateway au dashboard métier",
            "Basée sur ChirpStack v4 (open-source, mature) + couche SaaS propriétaire",
            "Portail web style Orange Live Objects — expérience opérateur professionnelle",
            "Agent IA opérationnel (MCP) : provisioning et diagnostics en langage naturel",
            "Application eau Shengda intégrée : compteurs, vannes, fuites, jumeau numérique",
            "Déploiement Docker — VM Ubuntu, cloud ou on-premise au Sénégal",
            "Multi-tenant natif : un infrastructure, plusieurs clients isolés",
        ],
        subtitle="Sonatel · Orange IoT — LoRaWAN SaaS Platform",
    )

    add_architecture_slide(prs)

    add_section_slide(prs, "Fonctionnalités")

    add_content_slide(
        prs,
        "Gestion du réseau LoRaWAN",
        [
            "Applications, device profiles, devices OTAA (DevEUI, JoinEUI, AppKey)",
            "Gateways : CRUD, GPS, statut online/offline, last seen",
            "Région EU868 — Packet Forwarder UDP :1700, Basics Station TLS :3001",
            "Downlink : queue, fPort, confirmé/non confirmé, commandes vanne",
            "Événements device : uplink, join, ack — intégration ChirpStack temps réel",
            "RF Scan gateway : analyse spectre, détection interférences et pollueurs RF",
            "Provisioning automatique tenant ChirpStack à la création client",
        ],
    )

    add_content_slide(
        prs,
        "Data, analytics & intégrations",
        [
            "Ingestion MQTT temps réel → TimescaleDB (RSSI, SNR, DR, historique uplinks)",
            "Data Messages : historique payloads, décodage hex, filtres par device/application",
            "Decoders JavaScript ChirpStack — codec Shengda V1.6 intégré",
            "Analytics : trafic 24h, devices/gateways actifs, RSSI moyen, graphiques horaires",
            "Archivage payloads MinIO — téléchargement et traçabilité long terme",
            "Connecteurs sortants : HTTP webhook, MQTT/MQTTS, MCP bidirectionnel",
            "Règles IF/THEN : alertes automatiques sur RSSI, SNR, conditions métier",
        ],
    )

    add_two_column_slide(
        prs,
        "Supervision & intelligence",
        "NOC & anomalies",
        [
            "Tableau de bord NOC style Datadog (refresh 30s)",
            "KPI : uplinks, devices actifs, gateways, RSSI moyen",
            "Alertes règles et billing en temps réel",
            "Détection anomalies : devices silencieux",
            "Signal faible, SNR dégradé, pics de trafic",
            "Dashboards personnalisables par sélection de devices",
        ],
        "Agent IA (MCP)",
        [
            "Chat en français — console web intégrée",
            "« Liste les gateways », « Crée un device », « Diagnostique le réseau »",
            "~25 outils MCP : ChirpStack + compteurs eau",
            "LLM local Ollama (CPU) ou OpenAI — mode hybride",
            "Intégration Cursor IDE pour les équipes techniques",
            "Réduction drastique du temps d'onboarding et du support N1",
        ],
    )

    add_content_slide(
        prs,
        "Administration, sécurité & multi-tenancy",
        [
            "IAM Keycloak — realm lorawan : platform-admin, tenant-admin, operator, viewer",
            "Isolation tenant : JWT tenant_id + mapping ChirpStack automatique",
            "API keys par tenant (lwp_*) — scopes read/write pour intégrations SI",
            "Gestion tenants : création, suspension, suppression, membres",
            "Billing intégré : metering uplinks/devices/gateways, plans Starter/Pro/Enterprise",
            "Stripe Checkout pour upgrade self-service (roadmap Phase 3)",
            "FUOTA : mise à jour firmware over-the-air (multicast EU868, Class C)",
        ],
    )

    add_section_slide(prs, "Applications métier — Eau intelligente")

    add_content_slide(
        prs,
        "Solution Shengda — Compteurs eau intelligents",
        [
            "Protocole Application Layer V1.6 — codec ChirpStack natif",
            "Télérelevé : index m³, impulsions, batterie (V), état vanne, alarmes",
            "Commandes downlink : ouvrir / fermer / débourrer vanne, télérelevé forcé",
            "Paramétrage : intervalle de rapport (600 s – 24 h), heure de début",
            "Payload applicatif : 22 octets — trame LoRaWAN ~35 octets (DR5, confirmed)",
            "Historique relevés décodés + suivi commandes (pending / sent / ack)",
            "Production : device 8254812510001415 — uplinks toutes les 4 h, index 275,99 m³",
        ],
    )

    add_two_column_slide(
        prs,
        "Jumeau numérique & détection de fuites",
        "Jumeau numérique réseau",
        [
            "Visualisation compteur + gateway LoRaWAN (topologie réseau)",
            "Mode Simulation vs Live réseau (données ChirpStack temps réel)",
            "Métriques : RSSI, SNR, DR, débit (m³/h), batterie, vanne",
            "Santé device : nominal, fuite probable, attaque magnétique",
            "Commandes vanne depuis le jumeau",
            "Lien direct vers alertes fuites",
        ],
        "Détection de fuites",
        [
            "device_flow_alarm — alarme débit compteur",
            "flow_with_valve_closed — débit vanne fermée",
            "high_continuous_flow — débit continu élevé",
            "night_flow — consommation nocturne anormale",
            "flow_overload — surcharge débit",
            "Résumé : actives, critiques, warnings — résolution manuelle",
        ],
    )

    add_section_slide(prs, "Modèle économique")

    add_content_slide(
        prs,
        "Différenciation vs LNS propriétaires",
        [
            "Coût maîtrisé : open-source ChirpStack + stack propriétaire — pas de licence par device/gateway",
            "Souveraineté : code et données hébergés au Sénégal, personnalisation illimitée",
            "Vertical eau intégré : Shengda, fuites, jumeau — pas de développement client séparé",
            "IA opérationnelle native : agent MCP — aucun concurrent LNS ne l'offre en standard",
            "Multi-tenant SaaS : onboarding client en minutes (vs semaines chez Actility/Netmore)",
            "Extensibilité : decoders JS, rules, webhooks, MQTT, MCP — intégration SI flexible",
            "Time-to-market : plateforme opérationnelle aujourd'hui — Phases 0 à 2.5 livrées",
        ],
    )

    add_compliance_slide(prs)

    add_two_column_slide(
        prs,
        "Offre commerciale (plans SaaS)",
        "Starter — 49 €/mois",
        [
            "50 devices · 5 gateways · 100K uplinks/mois",
            "Analytics, rules, NOC, agent IA",
            "Intégrations HTTP/MQTT",
            "Idéal : POC, pilotes utilities",
        ],
        "Pro / Enterprise",
        [
            "Pro — 199 €/mois : 500 devices, FUOTA, anomalies ML",
            "Enterprise — 999 €/mois : 10 000 devices, API keys, support prioritaire",
            "Facturation à l'usage : metering uplinks/devices/gateways",
            "Modèle scalable pour déploiement national SUEZ / New Deal",
        ],
    )

    add_section_slide(prs, "Roadmap & conclusion")

    add_content_slide(
        prs,
        "Roadmap produit",
        [
            "Phase 0 ✅ — ChirpStack + Platform API + Agent MCP + Console",
            "Phase 1 ✅ — Ingestion MQTT → TimescaleDB + Rule Engine + Analytics",
            "Phase 2 ✅ — IAM Keycloak, NOC, billing metering",
            "Phase 2.5 ✅ — Provisioning auto tenant, API keys, suspend/delete",
            "Phase 3 🔜 — Stripe, invitations email, anomalies ML avancées, FUOTA production",
            "Phase 4 — Déploiement national : redondance, HA, millions de devices",
            "Phase 5 — Nouvelles verticales : énergie (SENELEC), agriculture, smart city",
        ],
    )

    add_content_slide(
        prs,
        "Prochaines étapes",
        [
            "Pilote SUEZ : déploiement compteurs eau sur zone test (Dakar / région)",
            "Alignement programme New Deal : cadrage gouvernance plateforme IoT nationale",
            "Dimensionnement réseau national : gateways, airtime, duty cycle EU868",
            "Certification et SLA : engagement disponibilité plateforme production",
            "Formation équipes opérations Sonatel (NOC, agent IA, support client)",
            "Démonstration live : console, jumeau numérique, détection fuites, agent IA",
        ],
    )

    add_title_slide(
        prs,
        "Merci",
        "Plateforme IoT LoRaWAN — Souveraine · Scalable · Intelligente",
    )

    prs.save(OUTPUT)
    print(f"Présentation générée : {OUTPUT}")


if __name__ == "__main__":
    build()
